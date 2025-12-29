"""
Generate an editable PPTX version of the data pipeline figure for the paper.

Output:
  article/figures/data_pipeline.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _add_card(slide, x, y, w, h, title, body_lines, title_hex, fill_hex="#F7F9FC", border_hex="#2C3E50"):
    # Outer rounded rect
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(fill_hex)
    card.line.color.rgb = _rgb(border_hex)
    card.line.width = Pt(1.25)

    # Title bar
    title_h = int(h * 0.22)
    tb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, title_h)
    tb.fill.solid()
    tb.fill.fore_color.rgb = _rgb(title_hex)
    tb.line.fill.background()  # no border

    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(2)

    # Body text box (transparent, on top of card)
    tx = slide.shapes.add_textbox(x + Pt(10), y + title_h + Pt(6), w - Pt(16), h - title_h - Pt(10))
    ttf = tx.text_frame
    ttf.word_wrap = True
    ttf.clear()
    for i, line in enumerate(body_lines):
        p = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(12)
        p.font.color.rgb = _rgb("#1F2D3D")
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    return card


def _add_arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = _rgb("#34495E")
    conn.line.width = Pt(1.5)
    conn.line.end_arrowhead = True
    return conn


def main():
    out_dir = Path("article") / "figures"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "data_pipeline.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # widescreen
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Layout in inches
    col1_x, col2_x, col3_x, col4_x = Inches(0.4), Inches(3.7), Inches(7.0), Inches(10.3)
    w = Inches(2.9)
    h1 = Inches(1.55)
    r1_y = Inches(0.45)

    h2 = Inches(2.05)
    r2_y = Inches(2.45)

    h3 = Inches(1.55)
    r3_y = Inches(5.2)

    # Cards (titles/colors match PNG)
    _add_card(
        slide,
        col1_x,
        r1_y,
        w,
        h1,
        "Raw Câmara Open Data",
        [
            "Roll-call vote sessions",
            "Proposition metadata",
            "Authors / parties / legislature",
            "Government orientation fields",
        ],
        title_hex="#2E86AB",
    )
    _add_card(
        slide,
        col2_x,
        r1_y,
        w,
        h1,
        "Ingestion + Harmonization",
        ["Download / load sources", "Normalize identifiers", "Parse dates and outcomes", "Remove duplicates"],
        title_hex="#2E86AB",
    )
    _add_card(
        slide,
        col3_x,
        r1_y,
        w,
        h1,
        "Join & Filtering",
        ["Join sessions ↔ proposition info", "Join engineered feature tables", "Keep clear outcomes (0/1)", "Sort chronologically by date"],
        title_hex="#2E86AB",
    )
    _add_card(
        slide,
        col4_x,
        r1_y,
        w,
        h1,
        "Leakage-safe Protocol",
        ["Features use only past info", "No look-ahead / no test leakage", "Chronological evaluation"],
        title_hex="#2E86AB",
    )

    _add_card(
        slide,
        col2_x,
        r2_y,
        w,
        h2,
        "Feature Engineering (per session i)",
        [
            "Gov. orientation: resolve GOV. vs Governo",
            "Coalition size: num_authors_trunc, >10 flag",
            "Author popularity: past success (merged)",
            "Party popularity: last K sessions (K=5)",
            "HAR: past outcomes for same proposition",
            "Missing: popularity/party_pop=0, HAR=0.5",
        ],
        title_hex="#27AE60",
    )

    _add_card(
        slide,
        col4_x,
        r2_y,
        w,
        h2,
        "Modeling + Evaluation",
        [
            "Train: first 80% (time-ordered)",
            "Test: last 20% (future)",
            "Models: VOTE-RAP, baselines, VIOLA-style",
            "Metrics: AUROC + F1 (Rejected)",
            "Rejected threshold: maximize F1 for class 0",
        ],
        title_hex="#A23B72",
    )

    _add_card(
        slide,
        col1_x,
        r3_y,
        w,
        h3,
        "Artifacts (Saved Files)",
        [
            "vote_sessions_full.csv",
            "author_popularity.csv",
            "party_popularity_*.csv",
            "proposition_history_*.csv",
            "comparison_*.png / .csv",
        ],
        title_hex="#8E44AD",
    )

    # Arrows
    mid_y = r1_y + h1 / 2
    _add_arrow(slide, col1_x + w, mid_y, col2_x, mid_y)
    _add_arrow(slide, col2_x + w, mid_y, col3_x, mid_y)
    _add_arrow(slide, col3_x + w, mid_y, col4_x, mid_y)

    _add_arrow(slide, col3_x + w / 2, r1_y + h1, col2_x + w / 2, r2_y)
    _add_arrow(slide, col4_x + w / 2, r1_y + h1, col4_x + w / 2, r2_y)
    _add_arrow(slide, col2_x + w, r2_y + h2 / 2, col4_x, r2_y + h2 / 2)
    _add_arrow(slide, col2_x, r2_y + Inches(0.2), col1_x + w, r3_y)

    # Arrow labels
    lbl1 = slide.shapes.add_textbox(col2_x + Inches(2.2), Inches(1.85), Inches(3.4), Inches(0.4))
    tf = lbl1.text_frame
    tf.text = "compute features after merge (using only past sessions)"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = _rgb("#34495E")

    lbl2 = slide.shapes.add_textbox(col4_x + Inches(0.3), Inches(1.75), Inches(2.4), Inches(0.5))
    tf = lbl2.text_frame
    tf.text = "apply temporal split & evaluation"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = _rgb("#34495E")

    lbl3 = slide.shapes.add_textbox(Inches(1.95), Inches(4.35), Inches(1.5), Inches(0.4))
    tf = lbl3.text_frame
    tf.text = "outputs / caches"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = _rgb("#34495E")

    # Footer caption
    footer = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    ft = footer.text_frame
    ft.text = "Dataset construction and leakage-safe feature computation workflow (VOTE-RAP)."
    p = ft.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = _rgb("#34495E")

    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()


